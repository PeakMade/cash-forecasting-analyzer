"""
PowerPoint Generator Service
Creates professional PPTX slides from analysis recommendations
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from typing import Dict, Any
import logging
import os

logger = logging.getLogger(__name__)


class PowerPointGenerator:
    """Generate PowerPoint presentations from recommendation data"""
    
    def __init__(self):
        # PeakMade/Redpoint brand colors
        self.brand_blue = RGBColor(65, 105, 225)  # Royal Blue
        self.brand_orange = RGBColor(255, 140, 0)  # Dark Orange
        self.dark_gray = RGBColor(51, 51, 51)
        self.light_gray = RGBColor(128, 128, 128)
        
    def generate_presentation(self, recommendation: Dict[str, Any], output_path: str) -> str:
        """
        Generate PowerPoint presentation from recommendation data
        
        Args:
            recommendation: Complete recommendation dictionary from RecommendationEngine
            output_path: Path where PPTX file should be saved
            
        Returns:
            Path to generated PPTX file
        """
        logger.info(f"Generating PowerPoint presentation for {recommendation.get('property_name', 'Unknown')}")
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Title & Executive Decision
            self._add_title_slide(prs, recommendation)
            
            # Slide 2: Executive Summary (Bullets)
            self._add_executive_summary_slide(prs, recommendation)
            
            # Slide 3: Financial Metrics
            self._add_financial_metrics_slide(prs, recommendation)
            
            # Slide 4: Market & Economic Analysis
            self._add_market_analysis_slide(prs, recommendation)
            
            # Slide 5: Detailed Rationale
            self._add_detailed_rationale_slide(prs, recommendation)
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"PowerPoint presentation saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint: {str(e)}")
            raise
    
    def _add_title_slide(self, prs: Presentation, recommendation: Dict[str, Any]):
        """Slide 1: Title with property name and executive decision"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Property Name (Top)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = recommendation.get('property_name', 'Unknown Property')
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.dark_gray
        title_para.alignment = PP_ALIGN.CENTER
        
        # Analysis Period
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = f"{recommendation.get('analysis_month', 'Unknown')} Analysis"
        date_para = date_frame.paragraphs[0]
        date_para.font.size = Pt(18)
        date_para.font.color.rgb = self.light_gray
        date_para.alignment = PP_ALIGN.CENTER
        
        # Executive Decision Badge (Center)
        decision = recommendation.get('decision', 'DO_NOTHING')
        confidence = recommendation.get('confidence', 'MEDIUM')
        
        # Decision box background
        decision_colors = {
            'CONTRIBUTE': RGBColor(220, 53, 69),  # Red
            'DISTRIBUTE': RGBColor(40, 167, 69),  # Green
            'DO_NOTHING': RGBColor(23, 162, 184)  # Blue
        }
        
        decision_box = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(2.5), Inches(3.5), Inches(5), Inches(1.5)
        )
        decision_box.fill.solid()
        decision_box.fill.fore_color.rgb = decision_colors.get(decision, self.brand_blue)
        decision_box.line.color.rgb = decision_colors.get(decision, self.brand_blue)
        
        # Decision text
        decision_frame = decision_box.text_frame
        decision_frame.text = f"{decision.replace('_', ' ')}"
        decision_para = decision_frame.paragraphs[0]
        decision_para.font.size = Pt(36)
        decision_para.font.bold = True
        decision_para.font.color.rgb = RGBColor(255, 255, 255)
        decision_para.alignment = PP_ALIGN.CENTER
        
        # Confidence Level
        conf_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
        conf_frame = conf_box.text_frame
        conf_frame.text = f"Confidence Level: {confidence}"
        conf_para = conf_frame.paragraphs[0]
        conf_para.font.size = Pt(16)
        conf_para.font.color.rgb = self.light_gray
        conf_para.alignment = PP_ALIGN.CENTER
        
    def _add_executive_summary_slide(self, prs: Presentation, recommendation: Dict[str, Any]):
        """Slide 2: Executive Summary with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "📊 Executive Summary"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.brand_blue
        
        # Bullet points
        bullets_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
        text_frame = bullets_box.text_frame
        text_frame.word_wrap = True
        
        executive_summary = recommendation.get('executive_summary', [])
        for i, bullet in enumerate(executive_summary):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = f"✓ {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.dark_gray
            p.space_before = Pt(12)
            p.level = 0
    
    def _add_financial_metrics_slide(self, prs: Presentation, recommendation: Dict[str, Any]):
        """Slide 3: Key Financial Metrics"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "📈 Key Metrics & Context"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.brand_blue
        
        # Get detailed rationale
        detailed = recommendation.get('detailed_rationale', {})
        
        # Property section
        property_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(1.2))
        prop_frame = property_box.text_frame
        prop_frame.word_wrap = True
        
        property_name = recommendation.get('property_name', 'Unknown')
        analysis_month = recommendation.get('analysis_month', 'Unknown')
        
        p = prop_frame.paragraphs[0]
        p.text = f"Property: {property_name}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.dark_gray
        
        prop_frame.add_paragraph()
        p2 = prop_frame.paragraphs[1]
        p2.text = f"Analysis Period: {analysis_month}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = self.light_gray
        
        # Decision amount
        amount = recommendation.get('amount')
        if amount:
            prop_frame.add_paragraph()
            p3 = prop_frame.paragraphs[2]
            decision = recommendation.get('decision', '')
            action_text = 'Contribution' if 'CONTRIBUTE' in decision else 'Distribution'
            p3.text = f"Recommended {action_text}: ${amount:,.0f}"
            p3.font.size = Pt(16)
            p3.font.bold = True
            p3.font.color.rgb = self.brand_orange
        
        # Key metrics from detailed rationale
        metrics_text = self._format_detailed_rationale_summary(detailed)
        metrics_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.5), Inches(4))
        metrics_frame = metrics_box.text_frame
        metrics_frame.word_wrap = True
        metrics_frame.text = metrics_text
        
        for para in metrics_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.color.rgb = self.dark_gray
            para.space_after = Pt(6)
    
    def _add_market_analysis_slide(self, prs: Presentation, recommendation: Dict[str, Any]):
        """Slide 4: Market & Economic Analysis"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "🌍 Market & Economic Analysis"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.brand_blue
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        detailed = recommendation.get('detailed_rationale', {})
        market_text = detailed.get('economic_context', 'Market analysis pending')
        
        # Clean up the text - remove separator lines
        if market_text and '=' not in market_text[:50]:  # If no header separator
            clean_text = market_text
        else:
            lines = [line for line in market_text.split('\n') if line.strip() and not line.startswith('=')]
            clean_text = '\n'.join(lines[1:]) if len(lines) > 1 else market_text  # Skip title line
        
        p = content_frame.paragraphs[0]
        p.text = clean_text[:1500]  # Limit text to fit on slide
        p.font.size = Pt(12)
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
    
    def _add_detailed_rationale_slide(self, prs: Presentation, recommendation: Dict[str, Any]):
        """Slide 5: Detailed Rationale & Recommendations"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "💡 Decision Rationale & Risk Assessment"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.brand_blue
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        detailed = recommendation.get('detailed_rationale', {})
        
        # Combine decision rationale and risk assessment
        decision_text = detailed.get('decision_rationale', '')
        risk_text = detailed.get('risk_assessment', '')
        
        combined_text = ""
        if decision_text and decision_text != 'Detailed analysis pending':
            # Clean up separator lines
            lines = [line for line in decision_text.split('\n') if line.strip() and not line.startswith('=')]
            combined_text = '\n'.join(lines[1:]) if len(lines) > 1 else decision_text
        
        if risk_text and risk_text != 'Detailed analysis pending':
            if combined_text:
                combined_text += "\n\n"
            lines = [line for line in risk_text.split('\n') if line.strip() and not line.startswith('=')]
            combined_text += '\n'.join(lines[1:]) if len(lines) > 1 else risk_text
        
        if not combined_text:
            combined_text = "Comprehensive analysis complete. See executive summary for key findings."
        
        p = content_frame.paragraphs[0]
        p.text = combined_text[:1800]  # Limit to fit on slide
        p.font.size = Pt(11)
        p.font.color.rgb = self.dark_gray
        p.space_after = Pt(8)
    
    def _format_detailed_rationale_summary(self, detailed: Dict[str, Any]) -> str:
        """Format detailed rationale into summary text"""
        summary_parts = []
        
        if 'cash_forecast_analysis' in detailed:
            cash_text = detailed['cash_forecast_analysis']
            # Extract key lines
            lines = [line.strip() for line in cash_text.split('\n') if line.strip() and not line.startswith('=')]
            summary_parts.append('\n'.join(lines[:8]))  # First 8 meaningful lines
        
        return "\n".join(summary_parts) if summary_parts else "Analysis complete. See detailed rationale for full context."
